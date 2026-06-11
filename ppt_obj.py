from pptx import Presentation

pptx_path = r"temp.pptx"

prs = Presentation(pptx_path)

for slide_idx, slide in enumerate(prs.slides, start=1):
    print(f"\n===== 슬라이드 {slide_idx} =====")

    for shape_idx, shape in enumerate(slide.shapes, start=1):
        print(f"\n[객체 {shape_idx}]")
        print("이름:", shape.name)
        print("타입:", shape.shape_type)
        print("위치:", shape.left, shape.top)
        print("크기:", shape.width, shape.height)

        # 텍스트 객체
        if shape.has_text_frame:
            text = shape.text.strip()
            print("텍스트:", text)

        # 표 객체
        if shape.has_table:
            print("표 객체")
            table = shape.table
            print("행 수:", len(table.rows))
            print("열 수:", len(table.columns))

            for r, row in enumerate(table.rows):
                values = []
                for cell in row.cells:
                    values.append(cell.text.strip())
                print(f"  {r+1}행:", values)

        # 차트 객체
        if shape.has_chart:
            print("차트 객체")
            chart = shape.chart
            print("차트 타입:", chart.chart_type)

        # 그림 객체
        if shape.shape_type == 13:
            print("그림 객체")